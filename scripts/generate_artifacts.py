from __future__ import annotations

import textwrap
from pathlib import Path
import sys

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
OUT = ROOT / "outputs"
APP_SCREENSHOT = OUT / "scenesound_app_screenshot.png"
RESULT_SCREENSHOT = OUT / "scenesound_results_screenshot.png"


SLIDES = [
    ("SceneSound Studio Pro", ["Final Week 15 polished multimodal system", "Foundation Models for Speech, Music, Sound, and Multimodal AI", "CS 5588 - May 5, 2026"]),
    ("Problem Statement", ["Creators and educators need fast animal-care visuals, soundtracks, SFX, and narration for short explainers.", "Week 13 image prompts and Week 14 audio prompts were useful separately, but not a complete product.", "The final challenge needs a polished workflow with stronger engineering, evaluation, UX, business value, and responsible AI."]),
    ("Week 13 and 14 Summary", ["Week 13: AI Animal Care & Pet Health Visualization System using Stable Diffusion-style structured image prompts.", "Week 14: SceneSound Studio, a working text-to-music and sound-effects prototype for creator audio.", "Week 15 direction: combine image, music, SFX, narration, safety review, and evaluation into one multimodal studio."]),
    ("Final Enhancements", ["Multimodal pipeline: animal-care scenario -> image prompt -> soundtrack prompt -> SFX cue sheet -> narration.", "Improved prompt templates, audience controls, mood controls, duration controls, and safety checks.", "Gradio product demo, real evaluation metrics, business model, final report, and presentation-ready deck."]),
    ("Architecture", ["Input: animal, care condition, environment, mood, duration, and target audience.", "Processing: prompt normalization, image prompt generation, music prompt generation, SFX cue planning, narration, safety review.", "Output: production-ready creative brief plus evaluation metrics for prompt alignment, realism, diversity, and latency."]),
    ("Models Used", ["Week 13 target: Stable Diffusion / ControlNet for controlled animal-care image generation.", "Week 14 target: MusicGen or AudioLDM-style models for music beds and sound effects.", "Week 15 target: optional TTS or speech model for narration, plus a structured local demo for reliable presentation."]),
    ("Prompt and Input Design", ["Structured inputs keep the system controllable: subject, condition, setting, mood, duration, and audience.", "Prompt templates reduce vague outputs and produce safer educational animal-care media.", "Negative/safety framing avoids graphic medical content, copyrighted music imitation, and unauthorized voice cloning."]),
    ("How to Use the Demo", ["Step 1: enter the animal, care scenario, environment, soundtrack mood, duration, and audience.", "Step 2: click Generate Multimodal Plan to create the full production brief.", "Step 3: review each tab: Image Prompt, Music Prompt, SFX Cue Sheet, Narration, Responsible AI, and Evaluation Metrics."]),
    ("Screenshots", ["Input screen shows a productized Gradio-style workflow for multimodal scenario creation.", "Results screen shows image prompt, music prompt, SFX cues, narration, responsible-AI review, and metrics.", "The screenshots demonstrate how Week 13 image work and Week 14 sound work become one final system."]),
    ("Sample Output", ["Image prompt: educational veterinary scene for a golden retriever with mild dehydration warning signs.", "Music prompt: hopeful 45-second background bed with bright marimba, clean guitar, and soft strings.", "Final outputs also include SFX timing cues, narration text, safety review, prompt alignment, realism, diversity, and latency."]),
    ("Evaluation Results", ["Evaluation metrics are saved in outputs/evaluation_results.csv.", "Prompt alignment measures whether the generated plan includes subject, condition, mood, education, and care context.", "Realism, diversity, latency, safety flags, and SFX cue count show better output quality than the Week 14 audio-only baseline."]),
    ("Baseline vs Final Comparison", ["Week 14 baseline: audio prompt generator for music beds and sound effects.", "Week 15 final: multimodal production studio with image prompts, soundtrack prompts, SFX cues, narration, safety review, and metrics.", "The final version has stronger engineering, better outputs, deeper evaluation, and clearer product value."]),
    ("Business Value", ["Target users: veterinary educators, animal shelters, pet-care creators, instructors, and small marketing teams.", "Value: faster educational media production without hiring separate visual, music, and narration specialists.", "Pricing: free student tier, creator subscription, and clinic/shelter team tier with brand presets and export tools."]),
    ("Risks and Ethics", ["Avoid graphic or misleading animal-health visuals and label outputs as AI-generated educational media.", "Prevent copyright misuse by generating original music instead of imitating living artists or protected tracks.", "Prevent voice misuse by requiring permission for cloned voices and using neutral narration by default."]),
    ("Limitations", ["Offline demo generates prompts, cue sheets, narration, and metrics rather than heavy GPU image/audio files.", "Scores are proxy metrics; final production should include human ratings and domain expert review.", "Real model integration would require GPU/API access, dataset licensing, and stronger media validation."]),
    ("Future Work", ["Connect image prompts to Stable Diffusion XL or ControlNet and audio prompts to MusicGen or AudioLDM.", "Add waveform preview, generated image preview, timeline editing, batch export, and human rating UI.", "Add shelter/clinic templates, API deployment, and safer review workflows for educational publishing."]),
    ("Conclusion", ["SceneSound Studio Pro is significantly better than Week 14 because it turns an audio prototype into a full multimodal AI production workflow.", "It connects Week 13 image generation and Week 14 sound generation into a polished, evaluated, responsible, presentation-ready final system."]),
]


REPORT_SECTIONS = [
    ("Executive Summary", "SceneSound Studio Pro is a Week 15 final multimodal AI system that builds directly on Week 13 controlled animal-care image generation and Week 14 SceneSound audio generation. The final system creates image prompts, soundtrack prompts, SFX cue sheets, narration scripts, responsible-AI reviews, and evaluation metrics from one structured scenario."),
    ("Technical Approach", "The pipeline accepts animal, condition, environment, mood, duration, and audience inputs. It normalizes the inputs, creates image-generation prompts, creates music prompts, plans sound effects, writes a voice-ready narration, checks safety issues, and reports prompt alignment, realism, diversity, latency, safety flags, and cue counts."),
    ("Experiments and Evaluation", "The evaluation compares the Week 14 audio-only baseline against the Week 15 multimodal final system. Metrics include prompt alignment, multimodal quality, realism, diversity, latency, safety flags, and SFX cue coverage. This makes progress measurable instead of only visual."),
    ("Business Value", "Target users include veterinary educators, animal shelters, pet-care creators, instructors, and small marketing teams. The product reduces production time for educational media by combining visual direction, soundtrack planning, SFX, narration, and safety review in one workflow."),
    ("Responsible AI", "The project addresses graphic medical content, misleading veterinary advice, copyright risks in music generation, and voice cloning misuse. A production version should label AI outputs, avoid diagnosis claims, require licensing checks, and restrict cloned voice use without consent."),
    ("Lessons Learned", "The main lesson is that foundation-model projects become stronger when separate model capabilities are combined into a useful workflow. Week 15 improves engineering, UX, evaluation, business framing, and safety while keeping the same multimodal creative direction."),
    ("Conclusion", "SceneSound Studio Pro is significantly better than Week 14 because it turns a sound prompt prototype into a full presentation-ready multimodal AI production studio connected to both Week 13 and Week 14."),
]


THEME = {
    "navy": "172033",
    "ink": "25324A",
    "muted": "667085",
    "paper": "F7F9FC",
    "card": "FFFFFF",
    "blue": "2B7DE9",
    "teal": "0F9F6E",
    "gold": "D97706",
    "violet": "7C3AED",
    "rose": "D14368",
    "line": "D7DEE8",
}


def color(hex_value: str) -> RGBColor:
    hex_value = hex_value.lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def fill(shape, hex_value: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(hex_value)


def line(shape, hex_value: str = "D7DEE8", width: float = 1.0) -> None:
    shape.line.color.rgb = color(hex_value)
    shape.line.width = Pt(width)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int,
    hex_value: str = "25324A",
    bold: bool = False,
    align=PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Segoe UI"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color(hex_value)


def add_body(slide, bullets: list[str], x: float, y: float, w: float, h: float, size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {bullet}"
        paragraph.font.name = "Segoe UI"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color(THEME["ink"])
        paragraph.space_after = Pt(8)


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, accent: str) -> None:
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.04), Inches(y + 0.04), Inches(w), Inches(h))
    fill(shadow, "E9EEF6")
    shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(card, THEME["card"])
    line(card)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.08))
    fill(bar, accent)
    bar.line.fill.background()
    add_text(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.35, 16, accent, True)
    add_text(slide, body, x + 0.18, y + 0.62, w - 0.36, h - 0.72, 13, THEME["ink"])


def decorate(slide, prs, title: str, number: int) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color(THEME["paper"])
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.16))
    fill(top, THEME["blue"])
    top.line.fill.background()
    add_text(slide, title, 0.6, 0.42, 9.5, 0.6, 28, THEME["navy"], True)
    add_text(slide, "CS 5588 Week 15 Final Challenge | SceneSound Studio Pro", 0.65, 7.0, 9.0, 0.25, 9, THEME["muted"])
    add_text(slide, str(number), 12.25, 6.9, 0.45, 0.3, 11, THEME["muted"], True, PP_ALIGN.RIGHT)


def title_slide(slide, prs, bullets: list[str]) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color(THEME["navy"])
    fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)), THEME["navy"])
    for i, height in enumerate([1.0, 1.5, 0.75, 1.9, 1.25, 0.95, 1.65, 0.6]):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2 + i * 0.36), Inches(3.1 - height / 2), Inches(0.18), Inches(height))
        fill(bar, ["2B7DE9", "0F9F6E", "D97706", "7C3AED"][i % 4])
        bar.line.fill.background()
    add_text(slide, "SceneSound Studio Pro", 0.85, 1.05, 7.7, 0.9, 40, "FFFFFF", True)
    add_text(slide, "Animal-care visuals, soundtracks, SFX, narration, and safety review in one multimodal workflow.", 0.9, 2.05, 7.2, 0.7, 19, "DCE8F8")
    add_card(slide, "Final Week 15 polished system", bullets[1], 0.9, 3.15, 4.0, 1.45, THEME["blue"])
    add_card(slide, "Course", bullets[2], 5.15, 3.15, 3.25, 1.45, THEME["teal"])
    add_card(slide, "Demo Focus", "Image generation + music generation + narration planning + responsible AI.", 0.9, 4.9, 7.5, 1.15, THEME["violet"])
    add_text(slide, "01", 11.65, 6.75, 0.8, 0.35, 14, "DCE8F8", True, PP_ALIGN.RIGHT)


def cards_slide(slide, prs, title: str, bullets: list[str], number: int, accents: list[str]) -> None:
    decorate(slide, prs, title, number)
    widths = [3.85, 3.85, 3.85]
    for i, bullet in enumerate(bullets[:3]):
        add_card(slide, f"Point {i + 1}", bullet, 0.72 + i * 4.17, 1.65, widths[i], 3.25, accents[i % len(accents)])


def timeline_slide(slide, prs, title: str, bullets: list[str], number: int) -> None:
    decorate(slide, prs, title, number)
    labels = ["Week 13", "Week 14", "Week 15 Direction"]
    accents = [THEME["blue"], THEME["gold"], THEME["teal"]]
    for i, label in enumerate(labels):
        x = 0.82 + i * 4.12
        add_card(slide, label, bullets[i], x, 1.65, 3.62, 3.35, accents[i])
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 3.68), Inches(2.92), Inches(0.6), Inches(0.45))
            fill(arrow, "B8C7DD")
            arrow.line.fill.background()


def architecture_slide(slide, prs, title: str, bullets: list[str], number: int) -> None:
    decorate(slide, prs, title, number)
    steps = [
        ("Input", "Animal, condition, setting, mood, duration", THEME["blue"]),
        ("Prompt", "Image and soundtrack prompt templates", THEME["teal"]),
        ("Plan", "SFX cues, narration, and safety review", THEME["gold"]),
        ("Evaluate", "Alignment, realism, diversity, latency", THEME["violet"]),
    ]
    for i, (label, body, accent) in enumerate(steps):
        x = 0.65 + i * 3.1
        add_card(slide, label, body, x, 2.05, 2.55, 2.1, accent)
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.58), Inches(2.85), Inches(0.45), Inches(0.34))
            fill(arrow, "AAB7CA")
            arrow.line.fill.background()
    add_body(slide, bullets, 0.82, 5.05, 11.7, 0.95, 15)


def screenshot_slide(slide, prs, title: str, bullets: list[str], number: int) -> None:
    decorate(slide, prs, title, number)
    add_body(slide, bullets, 0.75, 1.18, 11.8, 0.85, 13)
    images = [(APP_SCREENSHOT, "Input workflow"), (RESULT_SCREENSHOT, "Generated multimodal plan")]
    for i, (path, caption) in enumerate(images):
        x = 0.75 + i * 6.2
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(5.75), Inches(3.7))
        fill(frame, "FFFFFF")
        line(frame, "CED8E8", 1.2)
        if path.exists():
            slide.shapes.add_picture(str(path), Inches(x + 0.15), Inches(2.48), width=Inches(5.45))
        add_text(slide, caption, x + 0.2, 6.13, 5.4, 0.32, 13, THEME["muted"], True, PP_ALIGN.CENTER)


def evaluation_slide(slide, prs, title: str, bullets: list[str], number: int) -> None:
    decorate(slide, prs, title, number)
    row = pd.read_csv(OUT / "evaluation_results.csv").iloc[0]
    metrics = [
        ("Prompt Alignment", f"{row['week15_prompt_alignment']:.0%}", THEME["teal"]),
        ("Multimodal Quality", f"{row['week15_multimodal_quality']:.0%}", THEME["blue"]),
        ("Diversity", f"{row['diversity_score']:.0%}", THEME["violet"]),
        ("Safety Flags", str(int(row["safety_flags_found"])), THEME["gold"]),
    ]
    for i, (label, value, accent) in enumerate(metrics):
        x = 0.75 + i * 3.05
        add_card(slide, label, value, x, 1.55, 2.55, 1.65, accent)
    add_body(slide, bullets, 1.05, 4.05, 11.2, 1.5, 18)


def comparison_slide(slide, prs, title: str, bullets: list[str], number: int) -> None:
    decorate(slide, prs, title, number)
    add_card(slide, "Baseline", bullets[0], 0.95, 1.65, 5.15, 3.4, THEME["rose"])
    add_card(slide, "Final Version", bullets[1], 7.15, 1.65, 5.15, 3.4, THEME["teal"])
    add_text(slide, bullets[2], 1.1, 5.65, 11.2, 0.55, 20, THEME["navy"], True, PP_ALIGN.CENTER)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for number, (title, bullets) in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if number == 1:
            title_slide(slide, prs, bullets)
        elif title == "Week 13 and 14 Summary":
            timeline_slide(slide, prs, title, bullets, number)
        elif title == "Architecture":
            architecture_slide(slide, prs, title, bullets, number)
        elif title == "Screenshots":
            screenshot_slide(slide, prs, title, bullets, number)
        elif title == "Evaluation Results":
            evaluation_slide(slide, prs, title, bullets, number)
        elif title == "Baseline vs Final Comparison":
            comparison_slide(slide, prs, title, bullets, number)
        elif title == "Conclusion":
            decorate(slide, prs, title, number)
            add_text(slide, bullets[0], 1.0, 1.8, 11.3, 1.2, 28, THEME["navy"], True, PP_ALIGN.CENTER)
            add_card(slide, "Final Takeaway", bullets[1], 2.2, 3.55, 8.9, 1.5, THEME["violet"])
        else:
            cards_slide(slide, prs, title, bullets, number, [THEME["blue"], THEME["teal"], THEME["gold"], THEME["violet"]])

    prs.save(OUT / "week15_scenesound_studio_pro.pptx")


def pdf_escape(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def wrap_lines(text: str, width: int = 92) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        lines.extend(textwrap.wrap(paragraph, width=width) or [""])
    return lines


def build_simple_pdf() -> None:
    lines = ["SceneSound Studio Pro - Final Report", "CS 5588 Week 15 Final Hands-On Challenge", ""]
    for title, body in REPORT_SECTIONS:
        lines.append(title)
        lines.extend(wrap_lines(body))
        lines.append("")

    pages = [lines[i : i + 42] for i in range(0, len(lines), 42)]
    objects: list[str] = []
    page_refs: list[int] = []

    def add_object(content: str) -> int:
        objects.append(content)
        return len(objects)

    font_obj = add_object("<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")

    for page in pages:
        stream_lines = ["BT", "/F1 11 Tf", "50 760 Td", "14 TL"]
        for line in page:
            stream_lines.append(f"({pdf_escape(line)}) Tj")
            stream_lines.append("T*")
        stream_lines.append("ET")
        stream = "\n".join(stream_lines)
        content_obj = add_object(f"<< /Length {len(stream.encode('latin-1', errors='replace'))} >>\nstream\n{stream}\nendstream")
        page_obj = add_object(f"<< /Type /Page /Parent 0 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 {font_obj} 0 R >> >> /Contents {content_obj} 0 R >>")
        page_refs.append(page_obj)

    kids = " ".join(f"{ref} 0 R" for ref in page_refs)
    pages_obj = add_object(f"<< /Type /Pages /Kids [{kids}] /Count {len(page_refs)} >>")
    catalog_obj = add_object(f"<< /Type /Catalog /Pages {pages_obj} 0 R >>")

    fixed_objects = []
    for idx, obj in enumerate(objects, start=1):
        if "/Parent 0 0 R" in obj:
            obj = obj.replace("/Parent 0 0 R", f"/Parent {pages_obj} 0 R")
        fixed_objects.append(obj)

    pdf_parts = ["%PDF-1.4\n"]
    offsets = [0]
    for idx, obj in enumerate(fixed_objects, start=1):
        offsets.append(sum(len(part.encode("latin-1", errors="replace")) for part in pdf_parts))
        pdf_parts.append(f"{idx} 0 obj\n{obj}\nendobj\n")
    xref_offset = sum(len(part.encode("latin-1", errors="replace")) for part in pdf_parts)
    pdf_parts.append(f"xref\n0 {len(fixed_objects) + 1}\n")
    pdf_parts.append("0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf_parts.append(f"{offset:010d} 00000 n \n")
    pdf_parts.append(f"trailer\n<< /Size {len(fixed_objects) + 1} /Root {catalog_obj} 0 R >>\nstartxref\n{xref_offset}\n%%EOF\n")
    (OUT / "week15_final_report.pdf").write_bytes("".join(pdf_parts).encode("latin-1", errors="replace"))


def main() -> None:
    OUT.mkdir(exist_ok=True)
    if not (OUT / "evaluation_results.csv").exists():
        from run_evaluation import main as run_eval

        run_eval()
    build_pptx()
    build_simple_pdf()
    print("Saved outputs/week15_scenesound_studio_pro.pptx")
    print("Saved outputs/week15_final_report.pdf")


if __name__ == "__main__":
    main()
